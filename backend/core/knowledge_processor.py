"""
Knowledge Processor — 다양한 파일 형식에서 텍스트 추출
지원: PDF, Excel, Word, PPT, Text, CSV, Markdown
"""
import logging
from pathlib import Path
from typing import List, Dict, Any

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    '.pdf': 'PDF 문서',
    '.xlsx': 'Excel 스프레드시트',
    '.xls': 'Excel 스프레드시트',
    '.docx': 'Word 문서',
    '.pptx': 'PowerPoint 프레젠테이션',
    '.txt': '텍스트 파일',
    '.csv': 'CSV 데이터',
    '.md': 'Markdown 문서',
}


def extract_text(file_path: str) -> List[Dict[str, Any]]:
    """
    파일에서 텍스트를 추출하여 청크 리스트로 반환.
    Returns: [{"text": str, "page": int, "source": str}, ...]
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == '.pdf':
        return _extract_pdf(file_path)
    elif ext in ('.xlsx', '.xls'):
        return _extract_excel(file_path)
    elif ext == '.docx':
        return _extract_word(file_path)
    elif ext == '.pptx':
        return _extract_ppt(file_path)
    elif ext in ('.txt', '.md'):
        return _extract_text(file_path)
    elif ext == '.csv':
        return _extract_csv(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")


def _chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """텍스트를 overlap이 있는 청크로 분할"""
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks


def _extract_pdf(file_path: str) -> List[Dict]:
    try:
        import pdfplumber
        chunks = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""

                # 이미지 속 텍스트 추출 (OCR):
                # - 텍스트가 매우 적은 페이지 (거의 이미지만)
                # - 또는 이미지가 포함된 페이지 (인포그래픽 등)
                try:
                    has_images = bool(page.images)
                except Exception:
                    has_images = False

                need_ocr = len(text.strip()) < 200 or has_images
                if need_ocr:
                    ocr_text = _ocr_pdf_page(file_path, page_num - 1)
                    if ocr_text:
                        # 중복 제거: OCR 텍스트에서 기존 텍스트와 겹치지 않는 부분만 추가
                        existing_words = set(text.lower().split())
                        ocr_words = ocr_text.split()
                        new_parts = []
                        for w in ocr_words:
                            if w.lower() not in existing_words:
                                new_parts.append(w)
                        if len(new_parts) > 10:  # 의미 있는 새 텍스트가 있을 때만
                            text = text + "\n[이미지 텍스트]\n" + ocr_text

                for chunk in _chunk_text(text):
                    chunks.append({"text": chunk, "page": page_num, "source": Path(file_path).name})
        return chunks
    except Exception as e:
        logger.error(f"PDF 추출 실패: {e}")
        return []


def _ocr_pdf_page(file_path: str, page_index: int) -> str:
    """PDF 페이지를 이미지로 변환 후 Claude Vision API로 OCR."""
    try:
        import fitz  # pymupdf
        import base64
        from backend.core.config import settings

        doc = fitz.open(file_path)
        page = doc[page_index]

        # 페이지를 PNG 이미지로 변환 (해상도 150 DPI)
        mat = fitz.Matrix(150 / 72, 150 / 72)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")
        doc.close()

        # Claude Vision API로 텍스트 추출
        import anthropic
        client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        img_b64 = base64.b64encode(img_bytes).decode("utf-8")

        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",  # 빠르고 저렴한 모델
            max_tokens=2000,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                    {"type": "text", "text": "이 이미지에서 모든 텍스트를 추출해주세요. 표, 도표, 인포그래픽 안의 텍스트도 모두 포함하세요. 텍스트만 반환하고 설명은 하지 마세요."},
                ],
            }],
        )
        return msg.content[0].text.strip() if msg.content else ""
    except Exception as e:
        logger.warning(f"OCR failed for page {page_index + 1}: {e}")
        return ""


def _extract_excel(file_path: str) -> List[Dict]:
    try:
        import pandas as pd
        chunks = []
        xl = pd.ExcelFile(file_path)
        for sheet_idx, sheet_name in enumerate(xl.sheet_names, 1):
            df = xl.parse(sheet_name, dtype=str).fillna("")
            # Convert sheet to readable text
            lines = [f"[시트: {sheet_name}]"]
            # Add headers
            lines.append(" | ".join(str(c) for c in df.columns))
            # Add rows
            for _, row in df.iterrows():
                row_text = " | ".join(str(v) for v in row.values if str(v).strip())
                if row_text.strip():
                    lines.append(row_text)
            sheet_text = "\n".join(lines)
            for chunk in _chunk_text(sheet_text):
                chunks.append({"text": chunk, "page": sheet_idx, "source": Path(file_path).name})
        return chunks
    except Exception as e:
        logger.error(f"Excel 추출 실패: {e}")
        return []


def _extract_word(file_path: str) -> List[Dict]:
    try:
        from docx import Document
        doc = Document(file_path)
        chunks = []
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text += "\n" + row_text
        for i, chunk in enumerate(_chunk_text(full_text)):
            chunks.append({"text": chunk, "page": i + 1, "source": Path(file_path).name})
        return chunks
    except Exception as e:
        logger.error(f"Word 추출 실패: {e}")
        return []


def _extract_ppt(file_path: str) -> List[Dict]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        chunks = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            for chunk in _chunk_text(slide_text):
                chunks.append({"text": chunk, "page": slide_num, "source": Path(file_path).name})
        return chunks
    except Exception as e:
        logger.error(f"PPT 추출 실패: {e}")
        return []


def _extract_text(file_path: str) -> List[Dict]:
    try:
        text = Path(file_path).read_text(encoding='utf-8', errors='replace')
        return [{"text": chunk, "page": i + 1, "source": Path(file_path).name}
                for i, chunk in enumerate(_chunk_text(text))]
    except Exception as e:
        logger.error(f"텍스트 추출 실패: {e}")
        return []


def _extract_csv(file_path: str) -> List[Dict]:
    try:
        import pandas as pd
        df = pd.read_csv(file_path, dtype=str).fillna("")
        lines = [" | ".join(str(c) for c in df.columns)]
        for _, row in df.iterrows():
            row_text = " | ".join(str(v) for v in row.values if str(v).strip())
            if row_text.strip():
                lines.append(row_text)
        text = "\n".join(lines)
        return [{"text": chunk, "page": i + 1, "source": Path(file_path).name}
                for i, chunk in enumerate(_chunk_text(text))]
    except Exception as e:
        logger.error(f"CSV 추출 실패: {e}")
        return []
